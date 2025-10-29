import os
import io
import json
import base64
import logging
import socket
from typing import List, Dict, Any, Optional, Union
from datetime import datetime
import hashlib
import openai
import tempfile
import subprocess
import shutil
import time
import urllib3
from urllib3.util.retry import Retry
from requests.adapters import HTTPAdapter
import requests
import ssl
from dotenv import load_dotenv

# 禁用SSL警告
urllib3.disable_warnings(urllib3.exceptions.InsecureRequestWarning)

# 加载环境变量
project_root = os.path.abspath(os.path.join(os.path.dirname(__file__), '..'))
load_dotenv(dotenv_path=os.path.join(project_root, '.env'))
load_dotenv(dotenv_path=os.path.join(project_root, '.env.local'))

# Pinecone兼容导入
try:
    from pinecone import Pinecone
    _PINECONE_CLIENT_MODE = 'new'
except Exception:
    Pinecone = None
    import pinecone as pinecone_legacy
    _PINECONE_CLIENT_MODE = 'legacy'

import openai
from pymongo import MongoClient
import requests
from PIL import Image
import pytesseract
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import cv2
import speech_recognition as sr
from moviepy.editor import VideoFileClip
import librosa
import soundfile as sf
import numpy as np
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity

# 配置日志
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class MultimediaProcessor:
    def __init__(self):
        # 初始化日志记录器
        self.logger = logging.getLogger(__name__)
        
        # 获取SiliconFlow API密钥
        self.api_key = os.getenv('SILICONFLOW_API_KEY')
        if not self.api_key:
            raise ValueError("SILICONFLOW_API_KEY 环境变量未设置")
        
        # 初始化SiliconFlow客户端
        self.client = openai.OpenAI(
            api_key=self.api_key,
            base_url="https://api.siliconflow.cn/v1"
        )
        
        # 初始化MongoDB
        self.mongo_client = MongoClient(os.getenv('MONGODB_URI', 'mongodb://localhost:27017/'))
        self.db = self.mongo_client['allpassagent']
        self.collection = self.db['multimedia_docs']
        self.chunks_collection = self.db['multimedia_chunks']
        
        # 初始化Pinecone (使用新版本API，参考document_processor的简洁写法)
        if _PINECONE_CLIENT_MODE == 'new':
            pc = Pinecone(api_key=os.getenv('PINECONE_API_KEY'))
            self.index = pc.Index(os.getenv('PINECONE_INDEX_NAME'))
        else:
            # 使用旧版本Pinecone
            import pinecone as pinecone_legacy
            pinecone_legacy.init(api_key=os.getenv('PINECONE_API_KEY'), environment=os.getenv('PINECONE_ENVIRONMENT'))
            self.index = pinecone_legacy.Index(os.getenv('PINECONE_INDEX_NAME'))

        # 支持的文件类型
        self.supported_types = {
            'ppt': ['.ppt', '.pptx'],
            'image': ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff'],
            'video': ['.mp4', '.avi', '.mov', '.wmv', '.flv', '.mkv'],
            'audio': ['.mp3', '.wav', '.flac', '.aac', '.ogg'],
            'document': ['.pdf', '.doc', '.docx', '.xls', '.xlsx', '.md', '.txt', '.html', '.htm', '.epub']
        }
        
        # 检查RAGAnything/MinerU是否可用
        self.raganything_available = self._check_raganything_available()





    def _safe_pinecone_query(self, query_vector, top_k=5, max_retries=3):
        """安全的Pinecone查询，带有智能重试和错误处理"""
        for attempt in range(max_retries):
            try:
                # 尝试查询
                results = self.index.query(
                    vector=query_vector,
                    top_k=top_k,
                    include_metadata=True
                )
                
                if results and hasattr(results, 'matches'):
                    self.logger.info(f"Pinecone查询成功，返回{len(results.matches)}个结果")
                    return results.matches
                else:
                    self.logger.warning("Pinecone查询返回空结果")
                    return []
                    
            except Exception as e:
                error_msg = str(e)
                self.logger.warning(f"遇到SSL EOF错误 (尝试 {attempt + 1}/{max_retries}): {error_msg}")
                
                if attempt < max_retries - 1:
                    # 在重试前等待
                    wait_time = (attempt + 1) * 2
                    self.logger.info(f"等待 {wait_time} 秒后重试...")
                    time.sleep(wait_time)
                    
                    # 重新初始化连接
                    try:
                        self._reinitialize_pinecone_connection()
                    except Exception as reinit_error:
                        self.logger.warning(f"重新初始化Pinecone连接失败: {reinit_error}")
                else:
                    self.logger.error(f"所有重试都失败了，最后错误: {error_msg}")
                    
        return []

    def _reinitialize_pinecone_connection(self):
        """重新初始化Pinecone连接"""
        try:
            # 重新初始化Pinecone客户端
            pinecone_api_key = os.getenv('PINECONE_API_KEY')
            if _PINECONE_CLIENT_MODE == 'new':
                pc = Pinecone(api_key=pinecone_api_key)
                self.index = pc.Index(os.getenv('PINECONE_INDEX_NAME'))
            else:
                import pinecone as pinecone_legacy
                pinecone_legacy.init(api_key=pinecone_api_key, environment=os.getenv('PINECONE_ENVIRONMENT'))
                self.index = pinecone_legacy.Index(os.getenv('PINECONE_INDEX_NAME'))
            
            self.logger.info("Pinecone连接重新初始化成功")
            return True
        except Exception as e:
            self.logger.error(f"重新初始化Pinecone连接失败: {e}")
            return False

    def process_multimedia_file(self, file_path: str, filename: str) -> Dict[str, Any]:
        """
        处理多媒体文件的主函数
        """
        try:
            file_ext = os.path.splitext(filename)[1].lower()
            file_type = self.get_file_type(file_ext)
            
            if not file_type:
                raise ValueError(f"不支持的文件类型: {file_ext}")
            
            logger.info(f"开始处理 {file_type} 文件: {filename}")
            
            # 创建文档记录
            doc_record = self.create_document_record(filename, file_path, file_type)
            doc_id = str(doc_record.inserted_id)
            
            # 根据文件类型选择处理方法
            if file_type == 'ppt':
                content_data = self.process_ppt(file_path)
            elif file_type == 'image':
                content_data = self.process_image(file_path)
            elif file_type == 'video':
                content_data = self.process_video(file_path)
            elif file_type == 'audio':
                content_data = self.process_audio(file_path)
            elif file_type == 'document':
                content_data = self.process_document_with_raganything(file_path)
            else:
                raise ValueError(f"未实现的文件类型处理: {file_type}")
            
            # 生成嵌入向量并存储
            self.store_multimedia_content(doc_id, filename, content_data, file_type)
            
            # 更新文档状态
            self.collection.update_one(
                {'_id': doc_record.inserted_id},
                {
                    '$set': {
                        'status': 'completed',
                        'processed_at': datetime.now(),
                        'content_count': len(content_data)
                    }
                }
            )
            
            logger.info(f"成功处理文件: {filename}")
            return {
                'success': True,
                'doc_id': doc_id,
                'content_count': len(content_data),
                'file_type': file_type
            }
            
        except Exception as e:
            logger.error(f"处理文件 {filename} 时出错: {str(e)}")
            if 'doc_record' in locals():
                self.collection.update_one(
                    {'_id': doc_record.inserted_id},
                    {'$set': {'status': 'failed', 'error': str(e)}}
                )
            return {'success': False, 'error': str(e)}

    def get_file_type(self, file_ext: str) -> Optional[str]:
        """根据文件扩展名确定文件类型"""
        file_ext = file_ext.lower()
        
        # 文档类型（使用RAGAnything处理）
        document_extensions = {
            '.pdf', '.doc', '.docx', '.xls', '.xlsx', 
            '.ppt', '.pptx', '.md', '.markdown', '.txt', 
            '.html', '.htm', '.epub', '.rtf', '.odt', 
            '.ods', '.odp', '.csv', '.tsv'
        }
        
        # 演示文稿类型（使用专门的PPT处理器）
        presentation_extensions = {'.ppt', '.pptx'}
        
        # 图像类型
        image_extensions = {
            '.jpg', '.jpeg', '.png', '.gif', '.bmp', 
            '.tiff', '.tif', '.webp', '.svg', '.ico'
        }
        
        # 视频类型
        video_extensions = {
            '.mp4', '.avi', '.mov', '.wmv', '.flv', 
            '.webm', '.mkv', '.m4v', '.3gp'
        }
        
        # 音频类型
        audio_extensions = {
            '.mp3', '.wav', '.flac', '.aac', '.ogg', 
            '.wma', '.m4a', '.opus'
        }
        
        # 优先级检查：先检查是否为PPT（使用专门的处理器）
        if file_ext in presentation_extensions:
            return 'ppt'
        elif file_ext in document_extensions:
            return 'document'
        elif file_ext in image_extensions:
            return 'image'
        elif file_ext in video_extensions:
            return 'video'
        elif file_ext in audio_extensions:
            return 'audio'
        else:
            return None

    def get_supported_file_types(self) -> Dict[str, List[str]]:
        """获取支持的文件类型详细信息"""
        return {
            'document': [
                '.pdf', '.doc', '.docx', '.xls', '.xlsx', 
                '.md', '.markdown', '.txt', '.html', '.htm', 
                '.epub', '.rtf', '.odt', '.ods', '.odp', 
                '.csv', '.tsv'
            ],
            'presentation': ['.ppt', '.pptx'],
            'image': [
                '.jpg', '.jpeg', '.png', '.gif', '.bmp', 
                '.tiff', '.tif', '.webp', '.svg', '.ico'
            ],
            'video': [
                '.mp4', '.avi', '.mov', '.wmv', '.flv', 
                '.webm', '.mkv', '.m4v', '.3gp'
            ],
            'audio': [
                '.mp3', '.wav', '.flac', '.aac', '.ogg', 
                '.wma', '.m4a', '.opus'
            ]
        }

    def validate_file_for_processing(self, file_path: str) -> Dict[str, Any]:
        """验证文件是否可以处理并返回详细信息"""
        if not os.path.exists(file_path):
            return {
                'valid': False,
                'error': '文件不存在',
                'file_type': None,
                'size': 0
            }
        
        file_size = os.path.getsize(file_path)
        max_size = 100 * 1024 * 1024  # 100MB限制
        
        if file_size > max_size:
            return {
                'valid': False,
                'error': f'文件过大 ({file_size / 1024 / 1024:.1f}MB)，最大支持100MB',
                'file_type': None,
                'size': file_size
            }
        
        filename = os.path.basename(file_path)
        file_ext = os.path.splitext(filename)[1]
        file_type = self.get_file_type(file_ext)
        
        if not file_type:
            return {
                'valid': False,
                'error': f'不支持的文件类型: {file_ext}',
                'file_type': None,
                'size': file_size
            }
        
        # 检查文档类型的解析器可用性
        if file_type == 'document':
            available_parsers = self._get_available_parsers()
            if not available_parsers:
                return {
                    'valid': False,
                    'error': '文档解析器不可用，请安装 raganything 或 mineru',
                    'file_type': file_type,
                    'size': file_size
                }
        
        return {
            'valid': True,
            'error': None,
            'file_type': file_type,
            'size': file_size,
            'filename': filename,
            'extension': file_ext,
            'available_parsers': self._get_available_parsers() if file_type == 'document' else None
        }

    def create_document_record(self, filename: str, file_path: str, file_type: str) -> Any:
        """
        创建文档记录
        """
        file_size = os.path.getsize(file_path) if os.path.exists(file_path) else 0
        
        doc_data = {
            'filename': filename,
            'original_name': filename,
            'file_type': file_type,
            'file_size': file_size,
            'uploaded_at': datetime.now(),
            'status': 'processing',
            'processed_at': None,
            'content_count': 0,
            'metadata': {}
        }
        
        return self.collection.insert_one(doc_data)

    def process_ppt(self, file_path: str) -> List[Dict[str, Any]]:
        """
        处理PPT文件，提取文本和图片信息
        """
        content_data = []
        
        try:
            prs = Presentation(file_path)
            
            for slide_idx, slide in enumerate(prs.slides):
                slide_content = {
                    'type': 'slide',
                    'slide_number': slide_idx + 1,
                    'text_content': '',
                    'images': [],
                    'shapes_info': []
                }
                
                # 提取文本内容
                text_parts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(shape.text.strip())
                        
                        # 记录形状信息
                        shape_info = {
                            'type': 'text',
                            'content': shape.text.strip(),
                            'shape_type': str(shape.shape_type)
                        }
                        slide_content['shapes_info'].append(shape_info)
                    
                    # 处理图片
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        try:
                            image_info = self.extract_image_from_shape(shape, slide_idx)
                            if image_info:
                                slide_content['images'].append(image_info)
                        except Exception as e:
                            logger.warning(f"提取幻灯片 {slide_idx + 1} 中的图片时出错: {str(e)}")
                
                slide_content['text_content'] = '\n'.join(text_parts)
                
                # 只有当幻灯片有内容时才添加
                if slide_content['text_content'] or slide_content['images']:
                    content_data.append(slide_content)
            
            logger.info(f"PPT处理完成，共提取 {len(content_data)} 张幻灯片")
            return content_data
            
        except Exception as e:
            logger.error(f"处理PPT文件时出错: {str(e)}")
            raise

    def extract_image_from_shape(self, shape, slide_idx: int) -> Optional[Dict[str, Any]]:
        """
        从PPT形状中提取图片信息
        """
        try:
            image = shape.image
            image_bytes = image.blob
            
            # 生成图片哈希
            image_hash = hashlib.md5(image_bytes).hexdigest()
            
            # 使用PIL处理图片
            pil_image = Image.open(io.BytesIO(image_bytes))
            
            # 提取图片中的文字（OCR）
            ocr_text = ""
            try:
                ocr_text = pytesseract.image_to_string(pil_image, lang='chi_sim+eng')
            except Exception as e:
                logger.warning(f"OCR处理失败: {str(e)}")
            
            return {
                'hash': image_hash,
                'format': image.content_type,
                'size': len(image_bytes),
                'dimensions': f"{pil_image.width}x{pil_image.height}",
                'ocr_text': ocr_text.strip(),
                'slide_number': slide_idx + 1
            }
            
        except Exception as e:
            logger.error(f"提取图片信息时出错: {str(e)}")
            return None

    def process_image(self, file_path: str) -> List[Dict[str, Any]]:
        """
        处理图片文件，提取文字和描述信息
        """
        content_data = []
        
        try:
            # 打开图片
            image = Image.open(file_path)
            
            # 基本信息
            image_info = {
                'type': 'image',
                'format': image.format,
                'size': image.size,
                'mode': image.mode,
                'ocr_text': '',
                'description': ''
            }
            
            # OCR文字识别
            try:
                ocr_text = pytesseract.image_to_string(image, lang='chi_sim+eng')
                image_info['ocr_text'] = ocr_text.strip()
            except Exception as e:
                logger.warning(f"OCR处理失败: {str(e)}")
            
            # 生成图片描述（这里可以集成图像识别API）
            image_info['description'] = self.generate_image_description(image)
            
            content_data.append(image_info)
            
            logger.info(f"图片处理完成，提取文字长度: {len(image_info['ocr_text'])}")
            return content_data
            
        except Exception as e:
            logger.error(f"处理图片文件时出错: {str(e)}")
            raise

    def generate_image_description(self, image: Image.Image) -> str:
        """
        生成图片描述（简单实现，可以扩展为调用视觉AI API）
        """
        try:
            # 这里可以集成如百度、腾讯、阿里等的图像识别API
            # 目前返回基本信息
            width, height = image.size
            aspect_ratio = width / height
            
            if aspect_ratio > 1.5:
                orientation = "横向"
            elif aspect_ratio < 0.67:
                orientation = "纵向"
            else:
                orientation = "方形"
            
            return f"这是一张{orientation}图片，尺寸为{width}x{height}像素"
            
        except Exception as e:
            logger.warning(f"生成图片描述时出错: {str(e)}")
            return "图片描述生成失败"

    def process_video(self, file_path: str) -> List[Dict[str, Any]]:
        """
        处理视频文件，提取音频转文字和关键帧
        """
        content_data = []
        
        try:
            # 使用moviepy处理视频
            video = VideoFileClip(file_path)
            duration = video.duration
            fps = video.fps
            
            video_info = {
                'type': 'video',
                'duration': duration,
                'fps': fps,
                'size': video.size,
                'audio_transcript': '',
                'keyframes': []
            }
            
            # 提取音频并转换为文字
            if video.audio:
                audio_transcript = self.extract_audio_transcript(video)
                video_info['audio_transcript'] = audio_transcript
            
            # 提取关键帧
            keyframes = self.extract_keyframes(video, max_frames=10)
            video_info['keyframes'] = keyframes
            
            video.close()
            content_data.append(video_info)
            
            logger.info(f"视频处理完成，时长: {duration:.2f}秒，提取关键帧: {len(keyframes)}个")
            return content_data
            
        except Exception as e:
            logger.error(f"处理视频文件时出错: {str(e)}")
            raise

    def extract_audio_transcript(self, video: VideoFileClip) -> str:
        """
        从视频中提取音频并转换为文字
        """
        try:
            # 提取音频
            audio = video.audio
            if audio is None:
                return "视频中没有音频轨道"
            
            # 统一转成16kHz PCM WAV，兼容多数ASR模型
            tmp = tempfile.NamedTemporaryFile(suffix='.wav', delete=False)
            temp_audio_path = tmp.name
            tmp.close()
            # 指定采样率与编解码器，避免默认参数造成兼容性问题
            audio.write_audiofile(
                temp_audio_path,
                fps=16000,
                codec='pcm_s16le',
                verbose=False,
                logger=None
            )

            # 基本健壮性检查，确保文件非空
            try:
                size = os.path.getsize(temp_audio_path)
                if size <= 44:  # 小于WAV头大小，视为异常文件
                    self.logger.warning(f"提取到的音频文件异常，大小={size} 字节")
                    return "音频提取失败：生成的WAV文件为空或损坏"
            except Exception as se:
                self.logger.warning(f"检查提取音频文件大小失败: {se}")

            # 使用语音识别，带重试机制
            return self._recognize_audio_with_retry(temp_audio_path)
            
        except Exception as e:
            self.logger.warning(f"音频转文字失败: {str(e)}")
            return f"音频转文字失败: {str(e)}"
        finally:
            # 清理临时文件
            try:
                if 'temp_audio_path' in locals() and os.path.exists(temp_audio_path):
                    os.remove(temp_audio_path)
            except Exception:
                pass

    def _recognize_audio_with_retry(self, audio_path: str, max_retries: int = 3) -> str:
        """
        带重试机制的语音识别，优先使用SiliconFlow API
        """
        for attempt in range(max_retries):
            try:
                # 首先尝试使用SiliconFlow API
                try:
                    transcript = self._recognize_with_siliconflow(audio_path)
                    if transcript and transcript.strip():
                        self.logger.info(f"SiliconFlow语音识别成功 (第{attempt+1}次尝试): {transcript[:50]}...")
                        return transcript
                except Exception as e:
                    self.logger.warning(f"SiliconFlow API识别失败 (第{attempt+1}次尝试): {str(e)}")
                
                # 如果SiliconFlow失败，回退到Google Speech API
                recognizer = sr.Recognizer()
                
                # 优化识别器参数
                recognizer.energy_threshold = 300
                recognizer.dynamic_energy_threshold = True
                recognizer.pause_threshold = 0.8
                recognizer.operation_timeout = 15  # 增加超时时间
                recognizer.phrase_time_limit = 30   # 增加短语时间限制
                
                with sr.AudioFile(audio_path) as source:
                    # 调整识别器参数以提高准确性
                    recognizer.adjust_for_ambient_noise(source, duration=0.5)
                    audio_data = recognizer.record(source)
                    
                    try:
                        # 使用Google语音识别（需要网络）
                        transcript = recognizer.recognize_google(
                            audio_data, 
                            language='zh-CN',
                            show_all=False
                        )
                        self.logger.info(f"Google语音识别成功 (第{attempt+1}次尝试): {transcript[:50]}...")
                        return transcript
                        
                    except sr.UnknownValueError:
                        self.logger.warning(f"无法识别音频内容 (第{attempt+1}次尝试)")
                        if attempt == max_retries - 1:
                            # 尝试离线识别引擎作为最终备用方案
                            return self._try_offline_recognition(audio_data)
                        continue
                        
                    except sr.RequestError as e:
                        error_msg = str(e).lower()
                        self.logger.warning(f"Google语音识别服务出错 (第{attempt+1}次尝试): {str(e)}")
                        
                        # 检查是否是网络连接问题
                        if any(keyword in error_msg for keyword in ['broken pipe', 'connection', 'timeout', 'network']):
                            if attempt < max_retries - 1:
                                self.logger.info(f"检测到网络问题，等待{2**attempt}秒后重试...")
                                time.sleep(2 ** attempt)  # 指数退避
                                continue
                        
                        # 尝试英文识别作为备用方案
                        if attempt == max_retries - 1:
                            try:
                                transcript_en = recognizer.recognize_google(
                                    audio_data, 
                                    language='en-US',
                                    show_all=False
                                )
                                self.logger.info(f"英文语音识别成功: {transcript_en[:50]}...")
                                return f"[英文识别] {transcript_en}"
                            except:
                                # 最终尝试离线识别
                                return self._try_offline_recognition(audio_data)
                        
            except Exception as e:
                self.logger.warning(f"语音识别过程异常 (第{attempt+1}次尝试): {str(e)}")
                if attempt < max_retries - 1:
                    time.sleep(1)
                    continue
                else:
                    # 最终尝试离线识别
                    try:
                        with sr.AudioFile(audio_path) as source:
                            recognizer = sr.Recognizer()
                            audio_data = recognizer.record(source)
                            return self._try_offline_recognition(audio_data)
                    except:
                        return f"语音识别失败: {str(e)}"
        
        return "语音识别重试次数已用完"

    def _recognize_with_siliconflow(self, audio_path: str) -> str:
        """
        使用SiliconFlow API进行语音识别
        """
        try:
            # 获取SiliconFlow API token
            siliconflow_token = os.getenv('SILICONFLOW_API_KEY')
            if not siliconflow_token:
                raise Exception("SILICONFLOW_API_KEY环境变量未设置")

            # 优先使用已初始化的OpenAI兼容客户端（base_url可在环境变量中配置）
            base_urls = [
                os.getenv('SILICONFLOW_BASE_URL', 'https://api.siliconflow.cn/v1'),
                'https://api.siliconflow.ai/v1',
            ]

            last_error = None
            for base_url in base_urls:
                try:
                    self.logger.info(f"使用SiliconFlow客户端进行转写: base_url={base_url}")
                    client = openai.OpenAI(api_key=siliconflow_token, base_url=base_url)
                    with open(audio_path, 'rb') as audio_file:
                        # OpenAI兼容API：audio.transcriptions.create
                        resp = client.audio.transcriptions.create(
                            model='TeleAI/TeleSpeechASR',
                            file=audio_file
                        )
                    # 兼容对象/字典两种返回结构
                    transcript = ''
                    if hasattr(resp, 'text'):
                        transcript = (resp.text or '').strip()
                    elif isinstance(resp, dict):
                        transcript = (resp.get('text') or '').strip()
                    else:
                        transcript = str(resp).strip()

                    if transcript:
                        return transcript
                    else:
                        raise Exception("API返回空的转录结果")

                except Exception as e:
                    self.logger.warning(f"SiliconFlow客户端调用失败(base_url={base_url}): {e}")
                    last_error = e
                    continue

            # 所有base_url都失败
            raise last_error if last_error else Exception("调用SiliconFlow API失败（未知错误）")
                    
        except Exception as e:
            self.logger.warning(f"SiliconFlow API调用失败: {str(e)}")
            raise e

    def _try_offline_recognition(self, audio_data) -> str:
        """
        尝试使用离线语音识别引擎
        """
        recognizer = sr.Recognizer()
        
        # 尝试使用Sphinx离线识别（如果可用）
        try:
            transcript = recognizer.recognize_sphinx(audio_data, language='zh-cn')
            if transcript and transcript.strip():
                self.logger.info(f"离线语音识别成功: {transcript[:50]}...")
                return f"[离线识别] {transcript}"
        except sr.UnknownValueError:
            self.logger.warning("离线语音识别无法识别音频内容")
        except sr.RequestError as e:
            self.logger.warning(f"离线语音识别不可用: {str(e)}")
        except Exception as e:
            self.logger.warning(f"离线语音识别异常: {str(e)}")
        
        # 如果离线识别也失败，返回音频基本信息
        return "语音识别服务暂时不可用，请检查网络连接或稍后重试"

    def extract_keyframes(self, video: VideoFileClip, max_frames: int = 10) -> List[Dict[str, Any]]:
        """
        从视频中提取关键帧
        """
        keyframes = []
        duration = video.duration
        
        try:
            # 均匀分布提取关键帧
            frame_times = np.linspace(0, duration - 1, max_frames)
            
            for i, time_point in enumerate(frame_times):
                try:
                    frame = video.get_frame(time_point)
                    
                    # 转换为PIL图片
                    pil_image = Image.fromarray(frame.astype('uint8'))
                    
                    # 生成缩略图
                    pil_image.thumbnail((200, 200))
                    
                    # 转换为base64
                    buffer = io.BytesIO()
                    pil_image.save(buffer, format='JPEG')
                    img_base64 = base64.b64encode(buffer.getvalue()).decode()
                    
                    keyframe_info = {
                        'frame_number': i + 1,
                        'timestamp': time_point,
                        'thumbnail_base64': img_base64,
                        'description': f"视频第{time_point:.1f}秒的关键帧"
                    }
                    
                    keyframes.append(keyframe_info)
                    
                except Exception as e:
                    logger.warning(f"提取第{i+1}个关键帧时出错: {str(e)}")
                    continue
            
            return keyframes
            
        except Exception as e:
            logger.error(f"提取关键帧时出错: {str(e)}")
            return []

    def process_audio(self, file_path: str) -> List[Dict[str, Any]]:
        """
        处理音频文件，转换为文字
        """
        content_data = []
        
        try:
            # 使用librosa加载音频
            y, sample_rate = librosa.load(file_path)
            duration = librosa.get_duration(y=y, sr=sample_rate)
            
            audio_info = {
                'type': 'audio',
                'duration': duration,
                'sample_rate': sample_rate,
                'transcript': ''
            }
            
            # 转换音频格式用于语音识别
            temp_wav_path = "/tmp/temp_audio_for_recognition.wav"
            sf.write(temp_wav_path, y, sample_rate)
            
            # 语音识别，使用重试机制
            transcript = self._recognize_audio_with_retry(temp_wav_path)
            audio_info['transcript'] = transcript
            
            content_data.append(audio_info)
            
            logger.info(f"音频处理完成，时长: {duration:.2f}秒")
            return content_data
            
        except Exception as e:
            logger.error(f"处理音频文件时出错: {str(e)}")
            raise
        finally:
            # 清理临时文件
            if os.path.exists("/tmp/temp_audio_for_recognition.wav"):
                os.remove("/tmp/temp_audio_for_recognition.wav")

    def store_multimedia_content(self, doc_id: str, filename: str, content_data: List[Dict[str, Any]], file_type: str):
        """
        存储多媒体内容到向量数据库
        """
        try:
            vectors_to_upsert = []
            chunks_to_store = []
            
            for idx, content in enumerate(content_data):
                # 构建文本内容用于向量化
                text_for_embedding = self.build_text_for_embedding(content, file_type)
                
                if text_for_embedding.strip():
                    # 生成嵌入向量
                    embedding = self.generate_embeddings(text_for_embedding)
                    
                    # 准备向量数据
                    vector_id = f"{doc_id}_{idx}"
                    
                    # 构建metadata，包含页码信息
                    metadata = {
                        'doc_id': doc_id,
                        'document_id': doc_id,  # 保持向后兼容
                        'filename': filename,
                        'file_type': 'multimedia',
                        'media_type': file_type,
                        'content_type': content.get('type', 'unknown'),
                        'chunk_index': idx,
                        'full_content': text_for_embedding[:1000],  # 限制长度
                        'content_summary': text_for_embedding[:200] + '...' if len(text_for_embedding) > 200 else text_for_embedding
                    }
                    
                    # 添加页码信息（如果存在）
                    if file_type == 'ppt' and content.get('slide_number'):
                        metadata['page'] = content['slide_number']
                        metadata['page_type'] = 'slide'
                    elif content.get('page_number'):
                        metadata['page'] = content['page_number']
                        metadata['page_type'] = 'page'
                    
                    vector_data = {
                        'id': vector_id,
                        'values': embedding,
                        'metadata': metadata
                    }
                    vectors_to_upsert.append(vector_data)
                    
                    # 准备MongoDB数据
                    chunk_data = {
                        'doc_id': doc_id,
                        'filename': filename,
                        'file_type': file_type,
                        'chunk_index': idx,
                        'content': content,
                        'text_content': text_for_embedding,
                        'created_at': datetime.now()
                    }
                    chunks_to_store.append(chunk_data)
            
            # 批量上传到Pinecone
            if vectors_to_upsert:
                logger.info(f"开始批量上传 {len(vectors_to_upsert)} 个向量到Pinecone")
                
                # 分批处理，每批100个向量
                batch_size = 100
                total_batches = (len(vectors_to_upsert) + batch_size - 1) // batch_size
                successful_uploads = 0
                failed_batches = []
                
                for batch_num in range(total_batches):
                    start_idx = batch_num * batch_size
                    end_idx = min(start_idx + batch_size, len(vectors_to_upsert))
                    batch = vectors_to_upsert[start_idx:end_idx]
                    
                    logger.info(f"上传批次 {batch_num + 1}/{total_batches} ({len(batch)} 个向量)")
                    
                    # 批次级别重试
                    max_retries = 5  # 增加重试次数
                    retry_delay = 2  # 初始延迟
                    
                    for attempt in range(max_retries):
                        try:
                            # 连接预检查
                            if attempt > 0:
                                logger.info(f"批次 {batch_num + 1} 第 {attempt + 1} 次尝试")
                                time.sleep(1)  # 短暂等待
                            
                            # 执行上传
                            self.index.upsert(vectors=batch)
                            successful_uploads += len(batch)
                            logger.info(f"批次 {batch_num + 1} 上传成功")
                            break
                            
                        except Exception as batch_error:
                            error_msg = str(batch_error)
                            logger.warning(f"批次 {batch_num + 1} 上传失败 (尝试 {attempt + 1}/{max_retries}): {error_msg}")
                            
                            # SSL错误特殊处理
                            if 'SSL' in error_msg or 'ssl' in error_msg.lower():
                                ssl_retry_delay = min(30, retry_delay * (3 ** attempt))  # SSL错误使用更长延迟
                                logger.info(f"SSL错误，延长等待时间到 {ssl_retry_delay} 秒")
                                time.sleep(ssl_retry_delay)
                                continue
                        
                            if attempt < max_retries - 1:
                                # 普通重试延迟
                                sleep_time = retry_delay * (2 ** attempt)
                                logger.info(f"等待 {sleep_time} 秒后重试...")
                                time.sleep(sleep_time)
                            else:
                                # 最终失败，记录到失败批次
                                failed_batches.append({
                                    'batch_num': batch_num,
                                    'batch_data': batch,
                                    'error': error_msg,
                                    'error_type': type(batch_error).__name__
                                })
                                logger.error(f"批次 {batch_num} 最终上传失败: {error_msg}")
                
                # 处理失败的批次
                if failed_batches:
                    logger.warning(f"部分批次上传失败: {len(failed_batches)}/{total_batches}")
                    logger.info(f"成功上传 {successful_uploads}/{len(vectors_to_upsert)} 个向量")
                    
                    # 降级方案：单个向量上传（仅对少量失败批次）
                    if len(failed_batches) <= 5:  # 增加降级方案的阈值
                        logger.info("尝试单个向量上传作为降级方案...")
                        recovered_vectors = 0
                        
                        for failed_batch in failed_batches:
                            logger.info(f"处理失败批次 {failed_batch['batch_num']}，包含 {len(failed_batch['batch_data'])} 个向量")
                            
                            for vector_idx, vector in enumerate(failed_batch['batch_data']):
                                try:
                                    # 单个向量也使用重试
                                    for single_attempt in range(3):
                                        try:
                                            self.index.upsert(vectors=[vector])
                                            recovered_vectors += 1
                                            break
                                        except Exception as single_retry_error:
                                            if single_attempt < 2:
                                                time.sleep(1 * (single_attempt + 1))
                                            else:
                                                raise single_retry_error
                                                
                                except Exception as single_error:
                                    logger.error(f"单个向量上传也失败 (批次 {failed_batch['batch_num']}, 向量 {vector_idx + 1}): {str(single_error)}")
                            
                        if recovered_vectors > 0:
                            successful_uploads += recovered_vectors
                            logger.info(f"通过单个上传恢复了 {recovered_vectors} 个向量")
                    else:
                        logger.warning(f"失败批次过多 ({len(failed_batches)})，跳过降级方案")
                
                # 最终结果评估
                if successful_uploads < len(vectors_to_upsert):
                    failure_rate = (len(vectors_to_upsert) - successful_uploads) / len(vectors_to_upsert) * 100
                    error_msg = f"Pinecone批量上传部分失败: 成功 {successful_uploads}/{len(vectors_to_upsert)} ({failure_rate:.1f}% 失败)"
                    
                    # 根据失败率决定是否抛出异常
                    if failure_rate > 70:  # 提高失败率阈值
                        logger.error(error_msg)
                        
                        # 针对不同错误类型的建议
                        if any('SSL' in fb.get('error', '') for fb in failed_batches):
                            logger.error("SSL连接问题诊断建议:")
                            logger.error("1. 检查网络连接稳定性")
                            logger.error("2. 确认Pinecone服务状态: https://status.pinecone.io/")
                            logger.error("3. 检查防火墙或代理设置")
                            logger.error("4. 验证系统时间是否正确")
                            logger.error("5. 尝试使用不同的网络环境")
                        
                        raise Exception(error_msg)
                    else:
                        logger.warning(f"{error_msg}，但失败率可接受，继续处理")
                else:
                    logger.info(f"所有批次上传成功: {successful_uploads}/{len(vectors_to_upsert)} 个向量")
                
            # 存储到MongoDB
            if chunks_to_store:
                self.chunks_collection.insert_many(chunks_to_store)
                logger.info(f"成功存储 {len(chunks_to_store)} 个内容块到MongoDB")
                
            logger.info(f"成功存储 {len(content_data)} 个内容块")
                
        except Exception as e:
            logger.error(f"存储多媒体内容时出错: {str(e)}")
            raise

    def build_text_for_embedding(self, content: Dict[str, Any], file_type: str) -> str:
        """
        构建用于向量化的文本内容
        根据内容类型和文件类型构建合适的文本表示
        """
        try:
            text_parts = []
            
            # 获取基础文本内容
            text_content = content.get('text_content', '')
            content_type = content.get('type', 'unknown')
            
            # 根据内容类型添加标识符
            if content_type == 'table':
                text_parts.append('[HTML表格]')
            elif content_type == 'equation':
                text_parts.append('[LaTeX公式]')
            elif content_type == 'image':
                text_parts.append('[图像内容]')
            elif content_type == 'text':
                text_parts.append('[文本内容]')
            elif content_type == 'slide':
                text_parts.append('[幻灯片内容]')
            elif content_type == 'video':
                text_parts.append('[视频内容]')
            elif content_type == 'audio':
                text_parts.append('[音频内容]')
            
            # 添加页码信息
            if content.get('page_number'):
                text_parts.append(f'[第{content["page_number"]}页]')
            elif content.get('slide_number'):
                text_parts.append(f'[第{content["slide_number"]}张幻灯片]')
            
            # 添加主要文本内容
            if text_content:
                text_parts.append(text_content)
            
            # 添加描述信息（如果存在）
            if content.get('description'):
                text_parts.append(f'描述: {content["description"]}')
            
            # 添加OCR文本（如果存在）
            if content.get('ocr_text'):
                text_parts.append(f'OCR识别: {content["ocr_text"]}')
            
            # 添加转录文本（如果存在）
            if content.get('transcript'):
                text_parts.append(f'语音转录: {content["transcript"]}')
            
            # 组合所有文本部分
            final_text = ' '.join(text_parts).strip()
            
            # 如果没有有效文本内容，返回基础信息
            if not final_text:
                final_text = f'[{file_type}文件] [{content_type}内容]'
            
            return final_text
            
        except Exception as e:
            logger.error(f"构建嵌入文本时出错: {str(e)}")
            # 返回基础信息作为后备
            return f'[{file_type}文件] [{content.get("type", "unknown")}内容]'

    def _batch_generate_embeddings(self, texts: List[str]) -> List[List[float]]:
        """
        批量生成嵌入向量，提高效率
        """
        embeddings = []
        batch_size = 20  # 根据API限制调整
        
        for i in range(0, len(texts), batch_size):
            batch_texts = texts[i:i + batch_size]
            
            try:
                # 批量调用嵌入API
                batch_embeddings = []
                for text in batch_texts:
                    embedding = self.generate_embeddings(text)
                    batch_embeddings.append(embedding)
                
                embeddings.extend(batch_embeddings)
                logger.info(f"已生成批次 {i//batch_size + 1}/{(len(texts)-1)//batch_size + 1} 的嵌入向量")
                
            except Exception as e:
                logger.error(f"批量生成嵌入向量失败: {str(e)}")
                # 单个重试
                for text in batch_texts:
                    try:
                        embedding = self.generate_embeddings(text)
                        embeddings.append(embedding)
                    except Exception as retry_e:
                        logger.error(f"单个嵌入向量生成失败: {str(retry_e)}")
                        # 使用零向量作为占位符
                        embeddings.append([0.0] * 1024)  # 根据模型维度调整
        
        return embeddings

    def generate_embeddings(self, text: str) -> List[float]:
        """
        生成文本嵌入向量
        """
        try:
            # 截断文本以适应token限制
            if len(text) > 400:
                text = text[:400]
            
            response = self.client.embeddings.create(
                model="BAAI/bge-large-zh-v1.5",
                input=text
            )
            return response.data[0].embedding
        except Exception as e:
            logger.error(f"生成嵌入向量时出错: {str(e)}")
            raise

    def search_multimedia_content(self, query: str, file_types: Optional[List[str]] = None, top_k: int = 5) -> List[Dict[str, Any]]:
        """
        搜索多媒体内容，使用增强的错误处理和连接管理
        """
        max_retries = 3
        base_delay = 0.5
        
        for attempt in range(max_retries):
            try:
                # 生成查询向量
                query_embedding = self.generate_embeddings(query)
                
                # 构建过滤条件
                filter_conditions = {}
                if file_types:
                    filter_conditions['file_type'] = {'$in': file_types}
                
                # 在每次重试前稍作等待，避免连续请求
                if attempt > 0:
                    wait_time = base_delay * (2 ** (attempt - 1))
                    logger.info(f"第 {attempt + 1} 次尝试查询，等待 {wait_time:.1f} 秒...")
                    time.sleep(wait_time)
                
                # 使用安全的Pinecone查询方法
                search_results_matches = self._safe_pinecone_query(
                    query_vector=query_embedding,
                    top_k=min(top_k, 5)  # 限制返回数量
                )
                
                results = []
                for match in search_results_matches:
                    result = {
                        'score': match.score,
                        'filename': match.metadata.get('filename', ''),
                        'file_type': match.metadata.get('media_type', match.metadata.get('file_type', '')),
                        'content_type': match.metadata.get('content_type', ''),
                        'content': match.metadata.get('full_content', ''),
                        'summary': match.metadata.get('content_summary', '')
                    }
                    results.append(result)
                
                logger.info(f"成功查询到 {len(results)} 条结果")
                return results
                
            except Exception as e:
                error_msg = str(e)
                
                # 检查是否是SSL EOF错误
                if 'SSL' in error_msg and ('EOF' in error_msg or 'UNEXPECTED_EOF_WHILE_READING' in error_msg):
                    if attempt < max_retries - 1:
                        logger.warning(f"遇到SSL EOF错误 (尝试 {attempt + 1}/{max_retries}): {error_msg}")
                        continue
                    else:
                        logger.error("SSL连接持续失败，已达到最大重试次数")
                        # 尝试重新初始化Pinecone连接
                        try:
                            logger.info("尝试重新初始化Pinecone连接...")
                            if _PINECONE_CLIENT_MODE == 'new' and Pinecone:
                                self.index = self.pc.Index(os.getenv('PINECONE_INDEX_NAME', 'allpassagent'))
                            else:
                                self.index = pinecone_legacy.Index(os.getenv('PINECONE_INDEX_NAME', 'allpassagent'))
                            logger.info("Pinecone连接重新初始化成功")
                        except Exception as reinit_e:
                            logger.error(f"重新初始化Pinecone连接失败: {reinit_e}")
                        return []
                        
                # 处理其他类型的错误
                elif 'timeout' in error_msg.lower() or 'connection' in error_msg.lower():
                    if attempt < max_retries - 1:
                        logger.warning(f"遇到连接/超时错误 (尝试 {attempt + 1}/{max_retries}): {error_msg}")
                        continue
                    else:
                        logger.error(f"连接错误重试失败: {error_msg}")
                        return []
                else:
                    logger.error(f"查询失败: {error_msg}")
                    if attempt < max_retries - 1:
                        continue
                    return []
        
        return []

    def get_supported_types(self) -> Dict[str, List[str]]:
        """
        获取支持的文件类型
        """
        return self.supported_types

    def _check_raganything_available(self) -> bool:
        """检查RAGAnything/MinerU是否可用"""
        try:
            # 优先检查MinerU CLI
            if shutil.which('mineru'):
                return True
            # 检查raganything Python包（添加超时和错误处理）
            try:
                import raganything
                return True
            except ImportError:
                self.logger.info("RAGAnything包未安装，跳过相关功能")
                return False
            except Exception as e:
                self.logger.warning(f"RAGAnything包导入时出现网络错误，跳过: {e}")
                return False
        except Exception as e:
            self.logger.warning(f"检查RAGAnything可用性时出错: {e}")
            return False

    def _get_available_parsers(self) -> List[str]:
        """获取可用的解析器列表"""
        parsers = []
        
        # 检查MinerU CLI
        if shutil.which('mineru'):
            parsers.append('mineru')
        
        # 检查RAGAnything Python包
        try:
            import raganything
            parsers.append('raganything')
        except ImportError:
            pass
            
        # 检查Docling（如果可用）
        try:
            import docling
            parsers.append('docling')
        except ImportError:
            pass
            
        return parsers

    def process_document_with_raganything(self, file_path: str, parser: str = 'auto') -> List[Dict[str, Any]]:
        """
        使用RAGAnything/MinerU解析通用文档（PDF、Office、HTML、Markdown、EPUB等）。
        
        Args:
            file_path: 文档文件路径
            parser: 解析器选择 ('auto', 'mineru', 'raganything', 'docling')
        
        Returns:
            标准化的内容块列表，每个块包含类型、页码、文本内容与位置等信息
        """
        content_data: List[Dict[str, Any]] = []
        
        if not self.raganything_available:
            self.logger.warning('RAGAnything/MinerU 未安装或不可用，跳过处理')
            return []
        
        available_parsers = self._get_available_parsers()
        if not available_parsers:
            raise RuntimeError('未找到可用的文档解析器')
        
        # 自动选择解析器
        if parser == 'auto':
            if 'mineru' in available_parsers:
                parser = 'mineru'
            elif 'raganything' in available_parsers:
                parser = 'raganything'
            elif 'docling' in available_parsers:
                parser = 'docling'
            else:
                raise RuntimeError('未找到可用的解析器')
        
        # 验证选择的解析器是否可用
        if parser not in available_parsers:
            raise RuntimeError(f'解析器 {parser} 不可用。可用解析器: {available_parsers}')
        
        logger.info(f"使用解析器 {parser} 处理文档: {file_path}")
        
        try:
            if parser == 'mineru':
                content_data = self._parse_with_mineru(file_path)
            elif parser == 'raganything':
                content_data = self._parse_with_raganything_api(file_path)
            elif parser == 'docling':
                content_data = self._parse_with_docling(file_path)
            else:
                raise ValueError(f"不支持的解析器: {parser}")
                
            logger.info(f"成功解析文档，获得 {len(content_data)} 个内容块")
            return content_data
            
        except Exception as e:
            logger.error(f"文档解析失败 (解析器: {parser}): {str(e)}")
            # 如果指定解析器失败，尝试其他可用解析器
            if parser != 'auto' and len(available_parsers) > 1:
                logger.info("尝试使用其他可用解析器...")
                for fallback_parser in available_parsers:
                    if fallback_parser != parser:
                        try:
                            return self.process_document_with_raganything(file_path, fallback_parser)
                        except Exception as fallback_e:
                            logger.warning(f"备用解析器 {fallback_parser} 也失败: {str(fallback_e)}")
                            continue
            raise

    def _parse_with_mineru(self, file_path: str) -> List[Dict[str, Any]]:
        """使用MinerU CLI解析文档"""
        content_data = []
        
        with tempfile.TemporaryDirectory() as tmpdir:
            cmd = ['mineru', '-p', file_path, '-o', tmpdir, '-m', 'auto']
            result = subprocess.run(cmd, capture_output=True, text=True)
            
            if result.returncode != 0:
                raise RuntimeError(f"MinerU解析失败: {result.stderr}")

            # 查找content_list.json
            content_json_path = None
            for root, _, files in os.walk(tmpdir):
                for f in files:
                    if f == 'content_list.json':
                        content_json_path = os.path.join(root, f)
                        break
                if content_json_path:
                    break
                    
            if not content_json_path:
                raise FileNotFoundError('未找到 MinerU 输出的 content_list.json')

            with open(content_json_path, 'r', encoding='utf-8') as jf:
                content_list = json.load(jf)

            # 标准化为本项目的内容结构
            for idx, item in enumerate(content_list):
                block_type = item.get('type', 'text')
                page_idx = item.get('page_idx')
                page_number = (page_idx + 1) if isinstance(page_idx, int) else None
                text_content = item.get('content') or ''
                bbox = item.get('bbox')

                # 仅保留可向量化的文本类内容
                if block_type in ['text', 'table', 'equation'] and text_content.strip():
                    content_data.append({
                        'type': block_type,
                        'page_number': page_number,
                        'text_content': text_content,
                        'bbox': bbox,
                        'source': 'mineru',
                        'index': idx
                    })
                    
        return content_data

    def _parse_with_raganything_api(self, file_path: str) -> List[Dict[str, Any]]:
        """使用RAGAnything Python API解析文档"""
        try:
            # 添加超时和错误处理
            try:
                from raganything import RAGAnything
                from raganything.core.modal_processors import ModalProcessors
            except ImportError as e:
                self.logger.warning(f"RAGAnything包未安装: {e}")
                return []
            except Exception as e:
                self.logger.warning(f"RAGAnything包导入时出现网络错误: {e}")
                return []
            
            content_data = []
            
            # 初始化模态处理器
            modal_processors = ModalProcessors()
            
            # 根据文件类型选择处理器
            file_ext = os.path.splitext(file_path)[1].lower()
            
            if file_ext == '.pdf':
                # PDF处理
                result = modal_processors.process_pdf(file_path)
            elif file_ext in ['.doc', '.docx']:
                # Word文档处理
                result = modal_processors.process_office_document(file_path)
            elif file_ext in ['.ppt', '.pptx']:
                # PowerPoint处理
                result = modal_processors.process_office_document(file_path)
            elif file_ext in ['.xls', '.xlsx']:
                # Excel处理
                result = modal_processors.process_office_document(file_path)
            elif file_ext in ['.html', '.htm']:
                # HTML处理
                result = modal_processors.process_html(file_path)
            elif file_ext in ['.md', '.markdown']:
                # Markdown处理
                result = modal_processors.process_markdown(file_path)
            elif file_ext == '.txt':
                # 纯文本处理
                result = modal_processors.process_text(file_path)
            else:
                raise ValueError(f"RAGAnything不支持的文件类型: {file_ext}")
            
            # 标准化处理结果
            if isinstance(result, dict) and 'content' in result:
                content_list = result['content']
            elif isinstance(result, list):
                content_list = result
            else:
                content_list = [result]
            
            for idx, item in enumerate(content_list):
                if isinstance(item, dict):
                    block_type = item.get('type', 'text')
                    page_number = item.get('page_number') or item.get('page_idx')
                    text_content = item.get('content') or item.get('text') or str(item)
                    bbox = item.get('bbox')
                else:
                    block_type = 'text'
                    page_number = None
                    text_content = str(item)
                    bbox = None
                
                if text_content and text_content.strip():
                    content_data.append({
                        'type': block_type,
                        'page_number': page_number,
                        'text_content': text_content,
                        'bbox': bbox,
                        'source': 'raganything',
                        'index': idx
                    })
                    
            return content_data
            
        except ImportError as e:
            raise RuntimeError(f"RAGAnything Python包未正确安装: {str(e)}")
        except Exception as e:
            raise RuntimeError(f"RAGAnything API调用失败: {str(e)}")

    def _parse_with_docling(self, file_path: str) -> List[Dict[str, Any]]:
        """使用Docling解析文档"""
        try:
            import docling
            from docling.document_converter import DocumentConverter
            
            content_data = []
            
            # 初始化文档转换器
            converter = DocumentConverter()
            
            # 转换文档
            result = converter.convert(file_path)
            
            # 处理转换结果
            if hasattr(result, 'document') and hasattr(result.document, 'body'):
                for idx, element in enumerate(result.document.body):
                    if hasattr(element, 'text') and element.text.strip():
                        content_data.append({
                            'type': getattr(element, 'type', 'text'),
                            'page_number': getattr(element, 'page_number', None),
                            'text_content': element.text,
                            'bbox': getattr(element, 'bbox', None),
                            'source': 'docling',
                            'index': idx
                        })
            
            return content_data
            
        except ImportError as e:
            raise RuntimeError(f"Docling包未正确安装: {str(e)}")
        except Exception as e:
            raise RuntimeError(f"Docling解析失败: {str(e)}")

if __name__ == "__main__":
    # 测试代码
    processor = MultimediaProcessor()
    print("多媒体处理器初始化成功")
    print("支持的文件类型:", processor.get_supported_types())